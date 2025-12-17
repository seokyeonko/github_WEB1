"""
[프로그램 설명서]
이 스크립트는 지정된 폴더(및 하위 폴더) 내의 다양한 문서 파일에서 텍스트를 추출하여
개별 텍스트 파일(.txt)로 저장하는 도구입니다.

[주요 기능]
1. 지원 파일 형식:
   - 한글 (.hwp) : hwp5txt.exe를 자동으로 탐지하여 텍스트 추출
   - 엑셀 (.xlsx, .xls) : 시트별 텍스트 추출
   - 파워포인트 (.pptx) : 슬라이드별 텍스트 추출
   - PDF (.pdf) : pdfplumber 및 PyPDF2를 병행하여 텍스트 추출
   - 워드 (.docx) : 문단 및 표 데이터 추출

2. 주요 특징:
   - 폴더 경로를 입력받아 재귀적으로 파일 탐색
   - tqdm 라이브러리를 이용한 진행률(Progress Bar) 표시
   - 추출된 텍스트는 원본 파일 위치의 'txt output' 폴더에 자동 저장
   - 인코딩 자동 감지 및 에러 처리 포함

[사용 방법]
1. 스크립트 실행
2. 콘솔에 텍스트를 추출할 대상 폴더 경로 입력
3. 진행 상황 확인 및 결과물 확인

[필요 모듈 설치]
pip install pandas python-pptx PyPDF2 python-docx pdfplumber chardet tqdm openpyxl
"""
import os
import subprocess
import pandas as pd
from pptx import Presentation
import PyPDF2
from docx import Document
import pdfplumber
import chardet
import warnings
import shutil
from tqdm import tqdm

warnings.filterwarnings('ignore')


# -------------------------------------------------------------------
#  HWP5TXT 자동 경로 탐지
# -------------------------------------------------------------------
def detect_hwp5txt():
    """
    1) venv/Scripts/hwp5txt.exe
    2) PATH에 등록된 hwp5txt.exe
    순으로 탐지하여 경로 반환
    """
    possible_names = ["hwp5txt.exe", "hwp5txt"]

    # 1) venv 자동 탐지
    venv_bin = os.path.dirname(os.sys.executable)
    for name in possible_names:
        path1 = os.path.join(venv_bin, name)
        if os.path.exists(path1):
            return path1

    # 2) PATH 탐색
    for name in possible_names:
        path2 = shutil.which(name)
        if path2:
            return path2

    return None


HWP5TXT_PATH = detect_hwp5txt()


# -------------------------------------------------------------------
#  HWP 파일 추출 (hwp5txt.exe 기반)
# -------------------------------------------------------------------
def extract_text_from_hwp(file_path):
    if not HWP5TXT_PATH:
        return "[HWP 오류] hwp5txt.exe를 찾을 수 없습니다. venv/Scripts 폴더를 확인하세요."

    try:
        result = subprocess.run(
            [HWP5TXT_PATH, file_path],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE
        )

        # 인코딩 자동 감지
        detect = chardet.detect(result.stdout)
        encoding = detect.get("encoding", "utf-8")

        text = result.stdout.decode(encoding, errors='ignore')

        if text.strip():
            return text
        else:
            err_text = result.stderr.decode('utf-8', errors='ignore')
            return f"[hwp5txt 출력 없음]\n{err_text}"

    except Exception as e:
        return f"[hwp5txt 실행 오류] {str(e)}"


# -------------------------------------------------------------------
#  기존 추출기 유지
# -------------------------------------------------------------------
def extract_text_from_excel(file_path):
    text = ""
    try:
        df = pd.read_excel(file_path, sheet_name=None, engine='openpyxl')
        for sheet_name, sheet_data in df.items():
            text += f"\n=== {sheet_name} ===\n"
            text += sheet_data.fillna('').to_string(index=False, max_colwidth=100)
            text += "\n\n"
    except Exception as e:
        text = f"[엑셀 오류] {str(e)}"
    return text


def extract_text_from_pptx(file_path):
    text = ""
    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides, 1):
            text += f"\n=== 슬라이드 {i} ===\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    cleaned_text = ' '.join(shape.text.split())
                    text += cleaned_text + "\n"
    except Exception as e:
        text = f"[PPTX 오류] {str(e)}"
    return text


def extract_text_from_pdf(file_path):
    text = ""
    try:
        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                text += f"\n=== 페이지 {i} ===\n"
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"

        # pdfplumber 결과가 없는 경우 PyPDF2로 재시도
        if not text.strip():
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for i, page in enumerate(reader.pages, 1):
                    text += f"\n=== 페이지 {i} ===\n"
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
    except Exception as e:
        text = f"[PDF 오류] {str(e)}"
    return text


def extract_text_from_word(file_path):
    text = ""
    try:
        doc = Document(file_path)
        for paragraph in doc.paragraphs:
            if paragraph.style.name.startswith('Heading'):
                text += f"\n=== {paragraph.text} ===\n"
            elif paragraph.text.strip():
                cleaned_text = ' '.join(paragraph.text.split())
                text += cleaned_text + "\n"

        # 표 처리
        for table in doc.tables:
            text += "\n=== 표 데이터 ===\n"
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells]
                text += ' | '.join(row_text) + "\n"

    except Exception as e:
        text = f"[워드 오류] {str(e)}"
    return text


# -------------------------------------------------------------------
#  확장자별 추출기 연결
# -------------------------------------------------------------------
def extract_text(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    extractors = {
        '.xlsx': extract_text_from_excel,
        '.xls': extract_text_from_excel,
        '.pptx': extract_text_from_pptx,
        '.pdf': extract_text_from_pdf,
        '.docx': extract_text_from_word,
        '.hwp': extract_text_from_hwp
    }
    extractor = extractors.get(ext)
    if extractor:
        return extractor(file_path)
    else:
        return "[지원하지 않는 파일 형식입니다.]"


# -------------------------------------------------------------------
#  파일 탐색 및 저장
# -------------------------------------------------------------------
def find_files(folder, exts=None):
    if exts is None:
        exts = {'.xlsx', '.xls', '.pptx', '.pdf', '.docx', '.hwp'}
    file_list = []
    for root, _, files in os.walk(folder):
        for f in files:
            if os.path.splitext(f)[1].lower() in exts:
                file_list.append(os.path.join(root, f))
    return file_list


def save_text_to_output_folder(orig_file_path, text):
    parent_dir = os.path.dirname(orig_file_path)
    output_dir = os.path.join(parent_dir, "txt output")
    os.makedirs(output_dir, exist_ok=True)

    base_name = os.path.splitext(os.path.basename(orig_file_path))[0]
    output_file = os.path.join(output_dir, f"{base_name}_추출결과.txt")

    try:
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(text)
        print(f"→ [{base_name}] 텍스트 추출 완료: {output_file}")
    except Exception as e:
        print(f"[파일 저장 오류] {output_file}: {str(e)}")


# -------------------------------------------------------------------
#  메인 실행
# -------------------------------------------------------------------
def get_folder_path():
    folder_path = input("텍스트를 추출할 '폴더' 경로를 입력하세요: ").strip()
    if not folder_path or not os.path.isdir(folder_path):
        print(f"[오류] 입력하신 경로가 폴더가 아니거나 존재하지 않습니다: {folder_path}")
        return None
    return folder_path


def main():
    folder = get_folder_path()
    if not folder:
        return

    files = find_files(folder)
    if not files:
        print("해당 폴더 및 하위폴더에 지원되는 파일이 없습니다.")
        return

    print(f"\n총 {len(files)}개 파일에서 텍스트 추출을 시작합니다.\n")

    for file_path in tqdm(files, desc="파일 처리 중"):
        extracted_text = extract_text(file_path)
        save_text_to_output_folder(file_path, extracted_text)

    print("\n모든 처리가 완료되었습니다.")


if __name__ == "__main__":
    main()
